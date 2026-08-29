# -*- coding: utf-8 -*-
"""
PRAHARI - AI Kavach 5-slide submission deck generator.

Run:  python build_deck.py
Out:  PRAHARI_AI_Kavach.pptx  (16:9, native editable shapes, dark tactical theme)
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# --------------------------------------------------------------------------
# theme
# --------------------------------------------------------------------------
BG      = RGBColor(0x0A, 0x0E, 0x1A)   # deep navy-black ground
PANEL   = RGBColor(0x13, 0x1A, 0x2B)   # content panel
PANEL2  = RGBColor(0x18, 0x21, 0x3A)   # nested / chip fill
RULE    = RGBColor(0x1F, 0x2A, 0x44)   # hairline
TEXT    = RGBColor(0xE8, 0xED, 0xF7)   # primary
MUTED   = RGBColor(0x93, 0xA2, 0xC0)   # secondary
DIM     = RGBColor(0x5A, 0x6B, 0x8C)   # tertiary / footer
CYAN    = RGBColor(0x22, 0xD3, 0xEE)   # system, flow
AMBER   = RGBColor(0xF5, 0xA5, 0x24)   # threat, findings
GREEN   = RGBColor(0x34, 0xD3, 0x99)   # proof, pass
RED     = RGBColor(0xEF, 0x44, 0x44)   # fail marks
VIOLET  = RGBColor(0xA7, 0x8B, 0xFA)   # reasoning tier

SANS = "Segoe UI"
MONO = "Consolas"
SYM  = "Segoe UI Symbol"

W, H = 13.333, 7.5      # slide size, inches
M    = 0.45             # page margin
CW   = W - 2 * M        # content width = 12.433


# --------------------------------------------------------------------------
# primitives
# --------------------------------------------------------------------------
def rect(slide, x, y, w, h, fill=None, line=None, lw=0.75,
         shape=MSO_SHAPE.ROUNDED_RECTANGLE, radius=0.07):
    s = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    s.shadow.inherit = False
    if radius is not None and len(s.adjustments):
        s.adjustments[0] = radius
    if fill is None:
        s.fill.background()
    else:
        s.fill.solid()
        s.fill.fore_color.rgb = fill
    if line is None:
        s.line.fill.background()
    else:
        s.line.color.rgb = line
        s.line.width = Pt(lw)
    s.text_frame.word_wrap = True
    return s


def P(t=None, runs=None, sz=10, c=None, b=False, f=SANS, i=False,
      align=PP_ALIGN.LEFT, sb=0, sa=4, ls=1.12, hang=None):
    return {"t": t, "runs": runs, "sz": sz, "c": c or TEXT, "b": b, "f": f, "i": i,
            "align": align, "sb": sb, "sa": sa, "ls": ls, "hang": hang}


def R(t, sz=None, c=None, b=None, f=None, i=None):
    d = {"t": t}
    if sz is not None: d["sz"] = sz
    if c  is not None: d["c"]  = c
    if b  is not None: d["b"]  = b
    if f  is not None: d["f"]  = f
    if i  is not None: d["i"]  = i
    return d


def _hanging(p, amount):
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(Inches(amount)))
    pPr.set("indent", str(-Inches(amount)))


def write(tf, paras):
    tf.word_wrap = True
    for idx, spec in enumerate(paras):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = spec["align"]
        p.space_before = Pt(spec["sb"])
        p.space_after = Pt(spec["sa"])
        p.line_spacing = spec["ls"]
        if spec.get("hang"):
            _hanging(p, spec["hang"])
        runs = spec["runs"] or [{"t": spec["t"] or ""}]
        for rs in runs:
            r = p.add_run()
            r.text = rs.get("t", "")
            fo = r.font
            fo.size   = Pt(rs.get("sz", spec["sz"]))
            fo.bold   = rs.get("b",  spec["b"])
            fo.italic = rs.get("i",  spec["i"])
            fo.name   = rs.get("f",  spec["f"])
            fo.color.rgb = rs.get("c", spec["c"])


def text(slide, x, y, w, h, paras, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    write(tf, paras)
    return box


def in_shape(shape, paras, pad_l=0.18, pad_t=0.12, pad_r=0.18, pad_b=0.10,
             anchor=MSO_ANCHOR.TOP):
    tf = shape.text_frame
    tf.margin_left, tf.margin_right = Inches(pad_l), Inches(pad_r)
    tf.margin_top, tf.margin_bottom = Inches(pad_t), Inches(pad_b)
    tf.vertical_anchor = anchor
    write(tf, paras)


def chip_w(label, sz=8.5, mono=True):
    cw = sz * (0.552 if mono else 0.498) / 72.0
    return len(label) * cw + 0.26


def chip(slide, x, y, label, color, sz=8.5, h=0.24, mono=True, filled=False):
    w = chip_w(label, sz, mono)
    s = rect(slide, x, y, w, h,
             fill=(color if filled else PANEL2),
             line=(None if filled else color), lw=0.6, radius=0.5)
    in_shape(s, [P(label, sz=sz, c=(BG if filled else color), b=True,
                   f=(MONO if mono else SANS), sa=0, ls=1.0,
                   align=PP_ALIGN.CENTER)],
             pad_l=0.06, pad_r=0.06, pad_t=0.0, pad_b=0.0,
             anchor=MSO_ANCHOR.MIDDLE)
    return w


def flow_chips(slide, x, y, maxw, items, sz=8.5, h=0.24, gap=0.09, rowgap=0.06,
               dry=False):
    """items: (label, color[, filled]). Returns height used; dry=True measures only."""
    cx, cy = x, y
    for it in items:
        label, color = it[0], it[1]
        filled = it[2] if len(it) > 2 else False
        w = chip_w(label, sz)
        if cx > x and cx + w > x + maxw:
            cx, cy = x, cy + h + rowgap
        if not dry:
            chip(slide, cx, cy, label, color, sz=sz, h=h, filled=filled)
        cx += w + gap
    return (cy + h) - y


def metric_tile(slide, x, y, w, h, value, label, color, vsz=19, pad_t=0.13):
    s = rect(slide, x, y, w, h, fill=PANEL, line=RULE, lw=0.75, radius=0.10)
    in_shape(s, [
        P(value, sz=vsz, c=color, b=True, f=MONO, sa=3, ls=0.95,
          align=PP_ALIGN.CENTER),
        P(label, sz=8, c=MUTED, sa=0, ls=1.05, align=PP_ALIGN.CENTER),
    ], pad_t=pad_t, pad_b=0.06, anchor=MSO_ANCHOR.TOP)
    return s


def arrow_glyph(slide, x, y, w=0.22, h=0.30, color=None):
    text(slide, x, y, w, h,
         [P("→", sz=12, c=color or RULE, b=True, sa=0, align=PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)


# --------------------------------------------------------------------------
# slide chrome
# --------------------------------------------------------------------------
def base_slide(prs, num, accent=CYAN):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    text(slide, M, 7.06, 8.0, 0.26,
         [P("AI KAVACH  ·  PRAHARI  ·  AUTONOMOUS CYBER-REASONING SYSTEM",
            sz=7.5, c=DIM, b=True, sa=0)])
    text(slide, W - M - 4.0, 7.06, 4.0, 0.26,
         [P("%d / 5" % num, sz=7.5, c=DIM, b=True, sa=0, align=PP_ALIGN.RIGHT)])
    return slide


def titled_slide(prs, num, title, kicker, accent=CYAN):
    s = base_slide(prs, num, accent)
    text(s, M, 0.30, 8.6, 0.52, [P(title, sz=24, b=True, c=TEXT, sa=0, ls=1.0)])
    text(s, M, 0.30, CW, 0.52,
         [P(kicker, sz=8.5, b=True, c=MUTED, sa=0, align=PP_ALIGN.RIGHT)],
         anchor=MSO_ANCHOR.BOTTOM)
    rect(s, M, 0.94, CW, 0.022, fill=accent, shape=MSO_SHAPE.RECTANGLE, radius=None)
    return s


def section_panel(slide, x, y, w, h, num, heading, accent, bullets,
                  body_sz=9.5, head_sz=10.5, tag=None):
    s = rect(slide, x, y, w, h, fill=PANEL, line=RULE, lw=0.75, radius=0.06)
    rect(slide, x, y, 0.055, h, fill=accent, shape=MSO_SHAPE.RECTANGLE, radius=None)
    paras = [P(runs=[R(num + "   ", f=MONO, c=accent, b=True, sz=head_sz),
                     R(heading, c=accent, b=True, sz=head_sz)], sa=7, ls=1.0)]
    if tag:
        paras.append(P(tag, sz=7.5, c=DIM, b=True, sa=7, ls=1.0))
    for b in bullets:
        paras.append(P(runs=[R("•  ", c=accent, b=True, sz=body_sz),
                             R(b, c=MUTED, sz=body_sz)],
                       sa=6, ls=1.14, hang=0.155))
    in_shape(s, paras, pad_l=0.20, pad_r=0.16, pad_t=0.16, pad_b=0.12)
    return s


# ==========================================================================
# SLIDE 1 - Introduction, Ideation & Brief Description
# ==========================================================================
def slide1(prs):
    s = base_slide(prs, 1)

    text(s, M, 0.26, 7.0, 0.68,
         [P("PRAHARI", sz=40, b=True, c=CYAN, sa=0, ls=0.95)])
    text(s, M, 0.92, 8.6, 0.26,
         [P("Proof-carrying Reasoning And Hardening for Autonomous Remediation of Intrusions",
            sz=10, c=MUTED, sa=0)])
    text(s, M, 0.26, CW, 0.34,
         [P("AI KAVACH  ·  CYBER-REASONING SYSTEM CHALLENGE",
            sz=8, c=DIM, b=True, sa=0, align=PP_ALIGN.RIGHT)])
    text(s, M, 0.62, CW, 0.34,
         [P("SLIDE 1  ·  INTRODUCTION, IDEATION & BRIEF DESCRIPTION",
            sz=8.5, c=MUTED, b=True, sa=0, align=PP_ALIGN.RIGHT)])
    rect(s, M, 1.24, CW, 0.022, fill=CYAN, shape=MSO_SHAPE.RECTANGLE, radius=None)

    # positioning statement
    pos = rect(s, M, 1.38, CW, 0.60, fill=PANEL, line=RULE, lw=0.75, radius=0.10)
    rect(s, M, 1.38, 0.055, 0.60, fill=CYAN, shape=MSO_SHAPE.RECTANGLE, radius=None)
    in_shape(pos, [P(runs=[
        R("An air-gap-deployable cyber-reasoning system that finds, patches, and ",
          sz=14, c=TEXT),
        R("proves", sz=14, c=CYAN, b=True),
        R(" — with no human in the loop.", sz=14, c=TEXT)], sa=0, ls=1.0)],
        pad_l=0.24, anchor=MSO_ANCHOR.MIDDLE)

    # three panels
    pw, gap = (CW - 0.60) / 3.0, 0.30
    py, ph = 2.12, 3.26
    xs = [M, M + pw + gap, M + 2 * (pw + gap)]

    section_panel(s, xs[0], py, pw, ph, "01", "THE PROBLEM", AMBER, [
        "Armed Forces web and API software — logistics portals, C2 services, "
        "internal REST APIs — is patched on a human cycle measured in weeks. "
        "Adversary exploitation windows are measured in hours.",
        "Commercial scanners find but never fix, bury analysts in false positives, "
        "and are cloud-only — a non-starter on a classified network.",
        "Every hour a known flaw stays live is an hour of exposed operational surface.",
        "The finale target is a simulated Armed Forces environment — the same shape of "
        "software, under the same constraint: no source code leaves the network.",
    ], body_sz=9.5, tag="WHY THE CURRENT PATCH CYCLE CANNOT HOLD")

    section_panel(s, xs[1], py, pw, ph, "02", "THE GAP", CYAN, [
        "Autonomous cyber-reasoning works on C/C++ because a segfault is a free "
        "ground-truth oracle. AddressSanitizer says “this is a real bug” "
        "— for free, with no ambiguity.",
        "Web and API services have no segfault. An SQL injection returns HTTP 200.",
        "So fuzzers are blind, findings are noisy guesses, and an LLM’s “fix” "
        "cannot be verified. This is the gap PRAHARI is built to close.",
        "Bolting an LLM onto a scanner inherits the same blindness — it only rewrites "
        "what the scanner already guessed at.",
    ], body_sz=9.5, tag="WHY EXISTING CRS DESIGNS DO NOT TRANSFER")

    section_panel(s, xs[2], py, pw, ph, "03", "THE SOLUTION", GREEN, [
        "PySan — a runtime taint sanitizer that gives Python services the oracle "
        "they never had: tainted input reaching a dangerous sink becomes a "
        "deterministic, replayable crash.",
        "Proof-carrying patches — no fix is accepted on model confidence. Each must "
        "clear four machine-checkable gates and ships with the evidence attached.",
        "Local-first reasoning. One laptop, fully offline, no cloud dependency.",
        "The result of a run is not an alert. It is a finding, a patch, and the "
        "evidence that the patch works — produced end to end without an operator.",
    ], body_sz=9.5, tag="SUPPLY THE MISSING ORACLE, THEN THE MISSING PROOF")

    # closed-loop chevron strip
    labels = [("FIND", AMBER), ("PROVE", AMBER), ("PATCH", CYAN),
              ("RE-PROVE", GREEN), ("REGRESS", GREEN)]
    cw_, cy = 2.60, 5.52
    step = (CW - cw_) / 4.0
    for i, (lab, col) in enumerate(labels):
        c = rect(s, M + i * step, cy, cw_, 0.60, fill=PANEL, line=col, lw=1.0,
                 shape=MSO_SHAPE.CHEVRON, radius=None)
        in_shape(c, [P(lab, sz=12, c=col, b=True, sa=0, ls=1.0,
                       align=PP_ALIGN.CENTER)],
                 pad_l=0.24, pad_r=0.10, pad_t=0, pad_b=0,
                 anchor=MSO_ANCHOR.MIDDLE)
    text(s, M, 6.18, CW, 0.26,
         [P("A closed autonomous loop — the system does not stop at a finding, "
            "and does not stop at a patch.",
            sz=9.5, c=MUTED, sa=0, align=PP_ALIGN.CENTER)])

    # team block
    tb = rect(s, M, 6.46, CW, 0.50, fill=PANEL, line=RULE, lw=0.75, radius=0.10)
    in_shape(tb, [P(runs=[
        R("[TEAM NAME]", sz=11, c=TEXT, b=True),
        R("        [Member 1]   ·   [Member 2]   ·   [Member 3]",
          sz=10, c=MUTED)], sa=0, ls=1.0)],
        pad_l=0.24, anchor=MSO_ANCHOR.MIDDLE)
    return s


# ==========================================================================
# SLIDE 2 - Detailed Methodology
# ==========================================================================
def stage_card(slide, x, y, w, h, num, title, body, accent,
               border=None, lw=0.75, tag=None, gates=None, out=None):
    s = rect(slide, x, y, w, h, fill=PANEL, line=border or RULE, lw=lw, radius=0.06)
    rect(slide, x, y, w, 0.045, fill=accent, shape=MSO_SHAPE.RECTANGLE, radius=None)
    paras = []
    if tag:
        paras.append(P(tag, sz=7, c=accent, b=True, f=MONO, sa=5, ls=1.0))
    paras.append(P(runs=[R(num + "  ", f=MONO, sz=11, c=accent, b=True),
                         R(title, sz=10, c=TEXT, b=True)], sa=6, ls=1.05))
    paras.append(P(body, sz=9, c=MUTED, sa=0, ls=1.16))
    in_shape(s, paras, pad_l=0.17, pad_r=0.15, pad_t=0.16, pad_b=0.10)
    if gates:
        gy = y + h - 0.06 - (0.22 * len(gates)) - (0.04 * (len(gates) - 1))
        for i, g in enumerate(gates):
            gs = rect(slide, x + 0.17, gy + i * 0.26, w - 0.34, 0.22,
                      fill=PANEL2, line=GREEN, lw=0.5, radius=0.5)
            in_shape(gs, [P(g, sz=8, c=GREEN, b=True, f=MONO, sa=0, ls=1.0)],
                     pad_l=0.12, pad_r=0.08, pad_t=0, pad_b=0,
                     anchor=MSO_ANCHOR.MIDDLE)
    if out:
        oy = y + h - 0.44
        rect(slide, x + 0.17, oy, w - 0.34, 0.012, fill=RULE,
             shape=MSO_SHAPE.RECTANGLE, radius=None)
        text(slide, x + 0.17, oy + 0.07, w - 0.34, 0.34,
             [P(runs=[R("OUT  ", sz=7.5, c=accent, b=True, f=MONO),
                      R(out, sz=8, c=MUTED)], sa=0, ls=1.12)])
    return s


def slide2(prs):
    s = titled_slide(prs, 2, "Detailed Methodology",
                     "SLIDE 2  ·  STEP-BY-STEP APPROACH, IMPLEMENTATION STRATEGY & WORKFLOW")

    cw_ = (CW - 3 * 0.22) / 4.0          # 2.943
    step = cw_ + 0.22
    xs = [M + i * step for i in range(4)]
    r1y, r2y, ch = 1.16, 3.40, 2.06

    stage_card(s, xs[0], r1y, cw_, ch, "01", "INGEST & MAP",
               "Clone the target. tree-sitter AST plus call graph. Routes lifted from "
               "OpenAPI / Swagger, Flask and FastAPI decorators, and Django URLconf.",
               CYAN, out="attack-surface inventory + call graph")
    stage_card(s, xs[1], r1y, cw_, ch, "02", "STATIC PRE-FOCUS",
               "Semgrep, CodeQL and custom taint rules rank every source-to-sink path. "
               "The LLM reads only the ranked slices, never the whole repository — "
               "the context economy that keeps PRAHARI lightweight.",
               CYAN, out="ranked candidate sink paths")
    stage_card(s, xs[2], r1y, cw_, ch, "03", "AUTO-HARNESS SYNTHESIS",
               "The model writes one fuzz driver per route, seeded from spec examples "
               "and recorded traffic. No human writes a harness — this is what lets "
               "PRAHARI be aimed at unfamiliar infrastructure.",
               CYAN, out="per-route fuzz drivers + seed corpus")
    stage_card(s, xs[3], r1y, cw_, ch, "04", "PySan TAINT ORACLE",
               "Import-time hooks on dangerous sinks: cursor.execute, os.system, "
               "subprocess, pickle.loads, eval, open, ORM raw. Taint propagates from "
               "every request object. Reaching a sink is a deterministic crash.",
               AMBER, border=AMBER, lw=1.5, tag="◆ CORE INNOVATION",
               out="replayable PoV + full taint path")

    for i in range(3):
        arrow_glyph(s, xs[i] + cw_, r1y + ch / 2 - 0.15, 0.22, 0.30)

    stage_card(s, xs[0], r2y, cw_, ch, "05", "TRIAGE & DEDUP",
               "PoVs are clustered by taint-path signature, mapped to CWE, and scored for "
               "severity by reachability and authentication requirement. Served entirely "
               "by the on-device 3B model.",
               VIOLET, out="deduplicated, CWE-tagged, ranked findings")
    stage_card(s, xs[1], r2y, cw_, ch, "06", "PATCH SYNTHESIS",
               "N candidate minimal diffs, each constrained by the recorded taint path and "
               "ranked locally. Escalates to the Claude API only when every local candidate "
               "fails validation.",
               VIOLET, out="candidate patch diffs")

    w7 = cw_ * 2 + 0.22
    stage_card(s, xs[2], r2y, w7, ch, "07", "PROOF HARNESS  —  FOUR GATES, ALL MUST PASS",
               "A patch is accepted only when all four gates pass. Output: a signed proof "
               "bundle plus a patch PR.", GREEN, border=GREEN, lw=1.5,
               tag="◆ THIS IS WHERE THE FIX IS PROVEN",
               gates=["GATE 1   the PoV replays and is now blocked",
                      "GATE 2   the existing test suite is 100% green",
                      "GATE 3   10k benign requests: pre/post responses identical",
                      "GATE 4   no new taint path introduced by the patch"])

    arrow_glyph(s, xs[0] + cw_, r2y + ch / 2 - 0.15, 0.22, 0.30)
    arrow_glyph(s, xs[1] + cw_, r2y + ch / 2 - 0.15, 0.22, 0.30)

    # feedback loop
    fy = r2y + ch + 0.12
    rect(s, xs[1] + 0.05, fy, 0.34, 0.28, fill=AMBER,
         shape=MSO_SHAPE.LEFT_ARROW, radius=None)
    fb = rect(s, xs[1] + 0.45, fy, (M + CW) - (xs[1] + 0.45), 0.28,
              fill=PANEL, line=AMBER, lw=0.75, radius=0.5)
    in_shape(fb, [P(runs=[
        R("FEEDBACK LOOP   ", sz=8, c=AMBER, b=True, f=MONO),
        R("any gate fails → regenerate the patch with the failing gate as evidence "
          "— bounded at K iterations, then escalated for human review",
          sz=8.5, c=MUTED)], sa=0, ls=1.0)],
        pad_l=0.16, pad_r=0.12, pad_t=0, pad_b=0, anchor=MSO_ANCHOR.MIDDLE)

    # implementation strategy - 36-hour build order
    sy, sh = 6.00, 0.98
    rect(s, M, sy, CW, sh, fill=PANEL, line=RULE, lw=0.75, radius=0.07)
    rect(s, M, sy, 0.055, sh, fill=CYAN, shape=MSO_SHAPE.RECTANGLE, radius=None)
    text(s, M + 0.22, sy + 0.11, CW - 0.44, 0.22,
         [P(runs=[R("IMPLEMENTATION STRATEGY", sz=9, c=CYAN, b=True),
                  R("     36-hour build order — the oracle is built first because "
                    "nothing downstream works without it", sz=8.5, c=DIM)],
            sa=0, ls=1.0)])
    phases = [
        ("H 00–06", "PySan sink hooks and taint propagation, validated against a "
                    "known-vulnerable Django app"),
        ("H 06–16", "Route mapping, harness synthesis, orchestrator state machine, "
                    "local model tier"),
        ("H 16–28", "Patch synthesis loop with all four gates and the feedback path "
                    "closed"),
        ("H 28–36", "Dashboard, proof-bundle signing, benchmark run against the "
                    "seeded corpus"),
    ]
    pw_ = (CW - 0.44 - 3 * 0.16) / 4.0
    for i, (lab, body) in enumerate(phases):
        px = M + 0.22 + i * (pw_ + 0.16)
        rect(s, px, sy + 0.40, 0.035, 0.44, fill=CYAN,
             shape=MSO_SHAPE.RECTANGLE, radius=None)
        text(s, px + 0.12, sy + 0.38, pw_ - 0.12, 0.48,
             [P(lab, sz=8, c=CYAN, b=True, f=MONO, sa=2, ls=1.0),
              P(body, sz=8, c=MUTED, sa=0, ls=1.12)])
    return s


# ==========================================================================
# SLIDE 3 - Technology Stack / Block Diagram
# ==========================================================================
def slide3(prs):
    s = titled_slide(prs, 3, "Technology Stack & System Architecture",
                     "SLIDE 3  ·  TECHNOLOGIES, FRAMEWORKS, BLOCK DIAGRAM & EQUIPMENT")

    main_w, side_x, side_w = 9.55, 10.25, 2.63
    bands = [
        ("TARGET\nLAYER", MUTED, [
            ("Flask", MUTED), ("FastAPI", MUTED), ("Django", MUTED),
            ("Node / Express", MUTED), ("Docker-containerised services", MUTED),
            ("OpenAPI 3 / Swagger", MUTED)]),
        ("SENSING\nLAYER", AMBER, [
            ("Semgrep", AMBER), ("CodeQL", AMBER), ("tree-sitter", AMBER),
            ("OpenAPI parser", AMBER), ("Atheris", AMBER), ("Hypothesis", AMBER),
            ("schemathesis", AMBER), ("coverage.py", AMBER),
            ("PySan taint runtime", AMBER, True)]),
        ("REASONING\nLAYER", VIOLET, [
            ("asyncio orchestrator (state machine)", VIOLET),
            ("TIER-1 LOCAL  Ollama + Qwen2.5-Coder 3B/7B", VIOLET, True),
            ("TIER-2 CLOUD  Claude API (escalation only)", VIOLET),
            ("FAISS RAG  CWE + CVE + codebase", VIOLET),
            ("token / CPU budget meter", VIOLET)]),
        ("ASSURANCE\nLAYER", GREEN, [
            ("pytest regression harness", GREEN),
            ("differential replay engine", GREEN),
            ("git apply / rollback", GREEN),
            ("ed25519 proof-bundle signer", GREEN, True),
            ("SARIF export", GREEN)]),
        ("OPS\nLAYER", CYAN, [
            ("FastAPI control plane", CYAN), ("HTMX dashboard", CYAN),
            ("SQLite findings store", CYAN), ("Docker Compose", CYAN),
            ("structured JSON audit log", CYAN),
            ("runs fully offline", CYAN, True)]),
    ]

    by, bh, bgap = 1.22, 1.02, 0.155
    lab_w = 1.52
    for i, (name, col, chips) in enumerate(bands):
        y = by + i * (bh + bgap)
        rect(s, M, y, main_w, bh, fill=PANEL, line=RULE, lw=0.75, radius=0.06)
        lb = rect(s, M, y, lab_w, bh, fill=PANEL2, line=None, radius=0.06)
        in_shape(lb, [P(name, sz=9, c=col, b=True, sa=0, ls=1.08,
                        align=PP_ALIGN.CENTER)],
                 pad_l=0.06, pad_r=0.06, pad_t=0, pad_b=0, anchor=MSO_ANCHOR.MIDDLE)
        rect(s, M, y, 0.05, bh, fill=col, shape=MSO_SHAPE.RECTANGLE, radius=None)
        cx = M + lab_w + 0.16
        cmaxw = main_w - lab_w - 0.32
        used = flow_chips(None, cx, 0, cmaxw, chips, dry=True)
        flow_chips(s, cx, y + (bh - used) / 2.0, cmaxw, chips)
        if i < len(bands) - 1:
            text(s, M + lab_w / 2 - 0.15, y + bh - 0.03, 0.30, bgap + 0.06,
                 [P("▼", sz=10, c=DIM, b=True, sa=0, align=PP_ALIGN.CENTER)],
                 anchor=MSO_ANCHOR.MIDDLE)

    # equipment sidebar
    side_h = 5 * bh + 4 * bgap
    sp = rect(s, side_x, by, side_w, side_h, fill=PANEL, line=CYAN, lw=1.0, radius=0.06)
    in_shape(sp, [
        P("EQUIPMENT & FOOTPRINT", sz=9.5, c=CYAN, b=True, sa=9, ls=1.0),
        P("No specialised hardware. No cloud. No licence-bound tooling.",
          sz=8.5, c=MUTED, sa=9, ls=1.15),
    ] + [P(runs=[R("•  ", c=CYAN, b=True, sz=9), R(t, c=TEXT, sz=9)],
           sa=5, ls=1.12, hang=0.15)
         for t in ["Commodity x86 laptop",
                   "16 GB RAM",
                   "8 GB GPU — optional, for the local model only",
                   "Everything in the stack is open source or already licensed"]],
        pad_l=0.18, pad_r=0.16, pad_t=0.18, pad_b=0.12)

    tx, tw_ = side_x + 0.18, side_w - 0.36
    rect(s, tx, by + 2.22, tw_, 0.014, fill=RULE,
         shape=MSO_SHAPE.RECTANGLE, radius=None)
    text(s, tx, by + 2.30, tw_, 0.20,
         [P("DESIGN TARGETS", sz=7.5, c=DIM, b=True, sa=0)])
    metric_tile(s, tx, by + 2.52, tw_, 0.74, "≥ 90%",
                "of decisions resolved on-device by the 3B local tier", CYAN,
                vsz=17, pad_t=0.10)
    metric_tile(s, tx, by + 3.32, tw_, 0.74, "< 6 GB",
                "peak RAM for a full autonomous run", GREEN,
                vsz=17, pad_t=0.10)
    metric_tile(s, tx, by + 4.12, tw_, 0.74, "< 90 s",
                "cold start to the first captured PoV", GREEN,
                vsz=17, pad_t=0.10)
    text(s, tx, by + 4.98, tw_, 0.70,
         [P("Deployable inside an air-gapped Armed Forces network, on hardware "
            "already in service.", sz=8.5, c=AMBER, b=True, sa=0, ls=1.15)])
    return s


# ==========================================================================
# SLIDE 4 - Salient Features & Novelty
# ==========================================================================
def slide4(prs):
    s = titled_slide(prs, 4, "Salient Features & Novelty",
                     "SLIDE 4  ·  KEY FEATURES, INNOVATION & UNIQUE SELLING PROPOSITION")

    hw = (CW - 0.28) / 2.0      # 6.076
    hy, hh = 1.18, 2.42

    # USP 1
    u1 = rect(s, M, hy, hw, hh, fill=PANEL, line=AMBER, lw=1.5, radius=0.08)
    rect(s, M, hy, hw, 0.05, fill=AMBER, shape=MSO_SHAPE.RECTANGLE, radius=None)
    in_shape(u1, [
        P(runs=[R("USP 1   ", sz=8, c=BG, b=True, f=MONO)], sa=0, ls=1.0),
        P(runs=[R("THE MISSING ORACLE — PySan", sz=13, c=AMBER, b=True)],
          sa=6, ls=1.0),
        P("Autonomous CRS works on C/C++ because AddressSanitizer turns a memory bug "
          "into a segfault — free ground truth. Web services have no such signal. "
          "PySan is the ASan-equivalent for Python: it propagates taint from every "
          "request into every dangerous sink, converting silent injection, path "
          "traversal, SSRF, unsafe deserialization and broken-authorization bugs into "
          "deterministic, replayable crashes.", sz=9.5, c=MUTED, sa=0, ls=1.15),
    ], pad_l=0.22, pad_r=0.20, pad_t=0.16, pad_b=0.10)
    # the USP-1 tag chip, drawn over the panel head
    chip(s, M + 0.22, hy + 0.14, "USP 1", AMBER, sz=8, h=0.22, filled=True)
    text(s, M + 0.22, hy + 1.96, hw - 0.44, 0.30,
         [P("No existing tool pairs runtime taint tracking with autonomous patching.",
            sz=9.5, c=TEXT, b=True, sa=0, ls=1.1)])
    ex = rect(s, M + 0.22, hy + 1.62, hw - 0.44, 0.28, fill=PANEL2, line=RULE,
              lw=0.6, radius=0.10)
    in_shape(ex, [P(runs=[
        R("tainted(request.args[‘q’])", sz=8.5, c=AMBER, f=MONO, b=True),
        R("  reaches  ", sz=8.5, c=DIM, f=MONO),
        R("cursor.execute()", sz=8.5, c=AMBER, f=MONO, b=True),
        R("   ⇒   PoV, CWE-89", sz=8.5, c=GREEN, f=MONO, b=True)], sa=0, ls=1.0)],
        pad_l=0.12, pad_r=0.08, pad_t=0, pad_b=0, anchor=MSO_ANCHOR.MIDDLE)

    # USP 2
    x2 = M + hw + 0.28
    u2 = rect(s, x2, hy, hw, hh, fill=PANEL, line=GREEN, lw=1.5, radius=0.08)
    rect(s, x2, hy, hw, 0.05, fill=GREEN, shape=MSO_SHAPE.RECTANGLE, radius=None)
    in_shape(u2, [
        P(runs=[R("USP 2   ", sz=8, c=BG, b=True, f=MONO)], sa=0, ls=1.0),
        P(runs=[R("PROOF-CARRYING PATCHES", sz=13, c=GREEN, b=True)], sa=6, ls=1.0),
        P("A patch is never accepted on model confidence. It must clear four "
          "machine-checkable gates, and it ships with the evidence attached:",
          sz=9.5, c=MUTED, sa=6, ls=1.15),
    ], pad_l=0.22, pad_r=0.20, pad_t=0.16, pad_b=0.10)
    chip(s, x2 + 0.22, hy + 0.14, "USP 2", GREEN, sz=8, h=0.22, filled=True)
    flow_chips(s, x2 + 0.22, hy + 1.22, hw - 0.44, [
        ("G1  PoV replay blocked", GREEN),
        ("G2  tests 100% green", GREEN),
        ("G3  10k-request differential fuzz identical", GREEN),
        ("G4  no new taint path", GREEN),
    ], sz=8, h=0.24)
    text(s, x2 + 0.22, hy + 1.86, hw - 0.44, 0.40,
         [P(runs=[
             R("The brief asks us to prove the fix holds. ", sz=9.5, c=TEXT, b=True),
             R("This answers it literally — and the bundle is re-checkable by a "
               "change-control board without trusting us.", sz=9.5, c=MUTED)],
            sa=0, ls=1.12)])

    # supporting features
    fy, fh = 3.82, 1.06
    fw = (CW - 3 * 0.20) / 4.0
    fstep = fw + 0.20
    feats = [
        ("ZERO-HARNESS AUTONOMY", CYAN,
         "Point it at a repository and a base URL. It writes its own fuzz drivers "
         "from the spec — nothing to configure for unfamiliar infrastructure."),
        ("AIR-GAP FIRST, LIGHTWEIGHT", CYAN,
         "About 90% of decisions are resolved by a 3B on-device model plus "
         "deterministic analysers. Cloud is an optional escalation, never a dependency."),
        ("EXPLAINABLE & AUDITABLE", CYAN,
         "Every finding carries its taint path, PoV, CWE mapping and signed proof "
         "bundle — built to survive defence change-control review."),
        ("GENERALISES BY DESIGN", CYAN,
         "The oracle sits at framework level, not application level. New target, "
         "no new rules, no retraining, no re-tuning."),
    ]
    for i, (t, col, body) in enumerate(feats):
        c = rect(s, M + i * fstep, fy, fw, fh, fill=PANEL, line=RULE, lw=0.75,
                 radius=0.07)
        rect(s, M + i * fstep, fy, fw, 0.04, fill=col, shape=MSO_SHAPE.RECTANGLE,
             radius=None)
        in_shape(c, [P(t, sz=9, c=col, b=True, sa=6, ls=1.05),
                     P(body, sz=8.5, c=MUTED, sa=0, ls=1.15)],
                 pad_l=0.16, pad_r=0.14, pad_t=0.15, pad_b=0.10)

    # comparison strip
    ty, rh = 5.18, 0.312
    lab_w = 3.40
    mw = (CW - lab_w) / 5.0
    cols = ["FINDS", "PATCHES", "PROVES", "WEB / PYTHON", "OFFLINE"]
    rows = [
        ("Scanners — Snyk, Semgrep, Burp", ["Y", "N", "N", "Y", "P"], False),
        ("LLM autofix — Copilot Autofix", ["Y", "Y", "N", "Y", "N"], False),
        ("AIxCC-class CRS — C/C++ only", ["Y", "Y", "Y", "N", "N"], False),
        ("PRAHARI", ["Y", "Y", "Y", "Y", "Y"], True),
    ]
    text(s, M, ty - 0.005, lab_w, rh,
         [P("WHERE PRAHARI SITS", sz=8.5, c=MUTED, b=True, sa=0)],
         anchor=MSO_ANCHOR.MIDDLE)
    for j, cname in enumerate(cols):
        text(s, M + lab_w + j * mw, ty - 0.005, mw, rh,
             [P(cname, sz=8, c=(GREEN if cname == "PROVES" else MUTED), b=True,
                sa=0, align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    for i, (label, marks, hero) in enumerate(rows):
        y = ty + rh + i * rh
        rect(s, M, y, CW, rh - 0.03,
             fill=(PANEL2 if hero else PANEL),
             line=(CYAN if hero else RULE), lw=(1.0 if hero else 0.5), radius=0.10)
        text(s, M + 0.18, y, lab_w - 0.20, rh - 0.03,
             [P(label, sz=9, c=(CYAN if hero else MUTED), b=hero, sa=0)],
             anchor=MSO_ANCHOR.MIDDLE)
        for j, mk in enumerate(marks):
            glyph, col = {"Y": ("✓", GREEN), "N": ("✗", RED),
                          "P": ("partial", AMBER)}[mk]
            text(s, M + lab_w + j * mw, y, mw, rh - 0.03,
                 [P(glyph, sz=(8 if mk == "P" else 12),
                    c=col, b=True, f=(SANS if mk == "P" else SYM),
                    sa=0, align=PP_ALIGN.CENTER)], anchor=MSO_ANCHOR.MIDDLE)
    text(s, M, 6.72, CW, 0.28,
         [P("PRAHARI is the only point in this space that is autonomous, verifiable, "
            "web-native and air-gap-capable at the same time.",
            sz=9, c=TEXT, b=True, sa=0, align=PP_ALIGN.CENTER)],
         anchor=MSO_ANCHOR.MIDDLE)
    return s


# ==========================================================================
# SLIDE 5 - Final Deliverables
# ==========================================================================
def slide5(prs):
    s = titled_slide(prs, 5, "Final Deliverables",
                     "SLIDE 5  ·  EXPECTED OUTCOMES, PERFORMANCE OBJECTIVES & PROOF-OF-CONCEPT")

    pw = (CW - 0.28) / 2.0
    py, ph = 1.20, 2.46

    d = rect(s, M, py, pw, ph, fill=PANEL, line=RULE, lw=0.75, radius=0.08)
    rect(s, M, py, 0.05, ph, fill=CYAN, shape=MSO_SHAPE.RECTANGLE, radius=None)
    items = [
        ("PRAHARI engine + CLI", "one-command Docker Compose run against any target "
                                 "repository and base URL"),
        ("PySan taint runtime", "released as a standalone, reusable Python library"),
        ("Proof-bundle spec + verifier", "anyone can re-check a patch without trusting "
                                         "the system that produced it"),
        ("Operator dashboard", "live findings, taint paths, patch diffs and gate results"),
        ("Benchmark report", "against a seeded-vulnerability corpus of Django/Flask apps "
                             "plus real CVEs"),
    ]
    paras = [P("DELIVERABLES", sz=10.5, c=CYAN, b=True, sa=8, ls=1.0)]
    for i, (name, desc) in enumerate(items):
        paras.append(P(runs=[R("%d  " % (i + 1), sz=9, c=CYAN, b=True, f=MONO),
                             R(name + "  ", sz=9, c=TEXT, b=True),
                             R("—  " + desc, sz=9, c=MUTED)],
                       sa=6, ls=1.14, hang=0.17))
    in_shape(d, paras, pad_l=0.22, pad_r=0.18, pad_t=0.16, pad_b=0.10)

    # POC / demo terminal
    x2 = M + pw + 0.28
    t = rect(s, x2, py, pw, ph, fill=PANEL, line=RULE, lw=0.75, radius=0.08)
    rect(s, x2, py, 0.05, ph, fill=AMBER, shape=MSO_SHAPE.RECTANGLE, radius=None)
    in_shape(t, [
        P("PROOF-OF-CONCEPT — LIVE DEMONSTRATION", sz=10.5, c=AMBER, b=True,
          sa=6, ls=1.0),
        P("A vulnerability is injected into a running target; the operator touches "
          "nothing after this line:", sz=9, c=MUTED, sa=0, ls=1.14),
    ], pad_l=0.22, pad_r=0.18, pad_t=0.16, pad_b=0.10)

    term = rect(s, x2 + 0.22, py + 0.86, pw - 0.44, 1.20, fill=BG, line=RULE,
                lw=0.75, radius=0.06)
    lines = [
        [R("$ prahari run --repo ./target --url http://127.0.0.1:8000",
           sz=8, c=TEXT, f=MONO, b=True)],
        [R("[map]     ", sz=8, c=DIM, f=MONO), R("routes discovered ............ 34",
                                                 sz=8, c=MUTED, f=MONO)],
        [R("[static]  ", sz=8, c=DIM, f=MONO), R("ranked sink paths ............ 12",
                                                 sz=8, c=MUTED, f=MONO)],
        [R("[pysan]   ", sz=8, c=DIM, f=MONO),
         R("PoV captured   CWE-89   GET /api/units?q=", sz=8, c=AMBER, f=MONO, b=True)],
        [R("[patch]   ", sz=8, c=DIM, f=MONO),
         R("candidate 2 of 4 accepted   (local 3B tier)", sz=8, c=MUTED, f=MONO)],
        [R("[gates]   ", sz=8, c=DIM, f=MONO),
         R("G1 blocked  G2 118/118  G3 10k ok  G4 clean", sz=8, c=GREEN, f=MONO, b=True)],
        [R("[bundle]  ", sz=8, c=DIM, f=MONO),
         R("proof a91f7c..  SIGNED  ->  patch PR opened", sz=8, c=GREEN, f=MONO, b=True)],
    ]
    in_shape(term, [P(runs=l, sa=1, ls=1.06) for l in lines],
             pad_l=0.14, pad_r=0.10, pad_t=0.10, pad_b=0.08)
    text(s, x2 + 0.22, py + 2.10, pw - 0.44, 0.26,
         [P("Illustrative output — every line above is produced by the pipeline "
            "on Slide 2.", sz=8, c=DIM, i=True, sa=0)])

    # performance objectives
    text(s, M, 3.78, CW, 0.26,
         [P(runs=[R("PERFORMANCE OBJECTIVES", sz=10, c=TEXT, b=True),
                  R("     targets for the 36-hour finale — not results already "
                    "measured", sz=8.5, c=AMBER, b=True)],
            sa=0)])
    tiles = [
        ("≥ 85%", "detection rate on the seeded corpus", CYAN),
        ("≤ 5%",  "false-positive rate on confirmed findings", CYAN),
        ("< 10 min", "mean time from finding to proven patch", CYAN),
        ("100%",     "regression suite pass before a patch ships", GREEN),
        ("< 6 GB",   "peak RAM for a full autonomous run", GREEN),
        ("≥ 90%", "of decisions resolved fully offline", GREEN),
    ]
    tw = (CW - 5 * 0.16) / 6.0
    tstep = tw + 0.16
    for i, (v, lab, col) in enumerate(tiles):
        metric_tile(s, M + i * tstep, 4.08, tw, 1.02, v, lab, col)

    # finale readiness
    fr = rect(s, M, 5.34, CW, 0.80, fill=PANEL, line=CYAN, lw=1.0, radius=0.08)
    rect(s, M, 5.34, 0.055, 0.80, fill=CYAN, shape=MSO_SHAPE.RECTANGLE, radius=None)
    in_shape(fr, [
        P("FINALE READINESS", sz=9, c=CYAN, b=True, sa=4, ls=1.0),
        P("PRAHARI is aimed at a target by repository path and base URL alone — no "
          "per-application rules, no retraining, no reconfiguration. It can be pointed at "
          "the simulated Indian Armed Forces environment on day one, and every result it "
          "produces arrives with the evidence a reviewer needs to accept or reject it.",
          sz=9, c=MUTED, sa=0, ls=1.14),
    ], pad_l=0.24, pad_r=0.20, pad_t=0.12, pad_b=0.08)

    # submission chips
    text(s, M, 6.30, 1.9, 0.30,
         [P("SUBMITTED AS", sz=8.5, c=DIM, b=True, sa=0)], anchor=MSO_ANCHOR.MIDDLE)
    flow_chips(s, M + 1.30, 6.32, CW - 1.30, [
        ("GitHub repository  source + Docker Compose", CYAN),
        ("Recorded 3-minute demonstration", CYAN),
        ("Benchmark & methodology report", CYAN),
    ], sz=8.5, h=0.26)
    return s


# --------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(W)
    prs.slide_height = Inches(H)
    slide1(prs)
    slide2(prs)
    slide3(prs)
    slide4(prs)
    slide5(prs)
    out = os.path.join(os.path.dirname(os.path.abspath(__file__)),
                       "PRAHARI_AI_Kavach.pptx")
    prs.save(out)
    print("wrote %s  (%d slides)" % (out, len(prs.slides._sldIdLst)))


if __name__ == "__main__":
    main()
